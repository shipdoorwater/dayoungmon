import os
import logging
from typing import Optional, List
from pathlib import Path

from PIL import Image
import PyPDF2
from docx import Document
from pptx import Presentation

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    logging.warning("Tesseract OCR is not available. Image text extraction will be disabled.")

class FileProcessor:
    """파일에서 텍스트를 추출하는 클래스"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_from_pdf,
            '.docx': self._extract_from_docx,
            '.pptx': self._extract_from_pptx,
            '.txt': self._extract_from_txt,
            '.png': self._extract_from_image,
            '.jpg': self._extract_from_image,
            '.jpeg': self._extract_from_image,
            '.bmp': self._extract_from_image,
            '.tiff': self._extract_from_image
        }
    
    def extract_text(self, file_path: str) -> Optional[str]:
        """파일에서 텍스트를 추출"""
        try:
            path = Path(file_path)
            if not path.exists():
                logging.error(f"파일을 찾을 수 없습니다: {file_path}")
                return None
            
            file_ext = path.suffix.lower()
            if file_ext not in self.supported_formats:
                logging.error(f"지원되지 않는 파일 형식입니다: {file_ext}")
                return None
            
            extractor = self.supported_formats[file_ext]
            return extractor(file_path)
            
        except Exception as e:
            logging.error(f"파일 처리 중 오류 발생: {str(e)}")
            return None
    
    def _extract_from_pdf(self, file_path: str) -> Optional[str]:
        """PDF 파일에서 텍스트 추출"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logging.error(f"PDF 처리 오류: {str(e)}")
            return None
    
    def _extract_from_docx(self, file_path: str) -> Optional[str]:
        """DOCX 파일에서 텍스트 추출"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logging.error(f"DOCX 처리 오류: {str(e)}")
            return None
    
    def _extract_from_pptx(self, file_path: str) -> Optional[str]:
        """PPTX 파일에서 텍스트 추출"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logging.error(f"PPTX 처리 오류: {str(e)}")
            return None
    
    def _extract_from_txt(self, file_path: str) -> Optional[str]:
        """TXT 파일에서 텍스트 추출"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='cp949') as file:
                    return file.read()
            except Exception as e:
                logging.error(f"TXT 파일 인코딩 오류: {str(e)}")
                return None
        except Exception as e:
            logging.error(f"TXT 처리 오류: {str(e)}")
            return None
    
    def _extract_from_image(self, file_path: str) -> Optional[str]:
        """이미지 파일에서 OCR로 텍스트 추출"""
        if not TESSERACT_AVAILABLE:
            return "이미지 텍스트 추출을 위해 Tesseract OCR 설치가 필요합니다."
        
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='kor+eng')
            return text.strip()
        except Exception as e:
            logging.error(f"이미지 OCR 처리 오류: {str(e)}")
            return None
    
    def get_supported_formats(self) -> List[str]:
        """지원되는 파일 형식 목록 반환"""
        return list(self.supported_formats.keys())

def preprocess_text(text: str) -> str:
    """텍스트 전처리 함수"""
    if not text:
        return ""
    
    # 연속된 공백을 하나로
    import re
    text = re.sub(r'\s+', ' ', text)
    
    # 앞뒤 공백 제거
    text = text.strip()
    
    return text